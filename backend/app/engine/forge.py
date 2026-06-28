"""The Forge (v2) — turns Howler's tagged blocks into real, downloadable files.

Pure rendering, NO model calls: it takes the brief (a title + prose blocks) and renders Markdown,
HTML, PDF, DOCX, plus XLSX (a sheet of findings), PPTX (a slide per finding), and PNG (a summary
card). Each comes back as bytes the Supervisor saves as an artifact. Every renderer is best-effort —
a failure in one format never sinks the others.
"""

from __future__ import annotations

import io
from xml.sax.saxutils import escape

import markdown as _markdown
from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

# format -> MIME type for downloads.
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME: dict[str, str] = {
    "md": "text/markdown; charset=utf-8",
    "html": "text/html; charset=utf-8",
    "pdf": "application/pdf",
    "docx": _DOCX_MIME,
    "xlsx": _XLSX_MIME,
    "pptx": _PPTX_MIME,
    "png": "image/png",
}


def _title_and_paras(blocks: list[dict], title: str = "") -> tuple[str, list[str]]:
    """Split blocks into a title + body paragraphs. A leading '# Heading' becomes the title."""
    paras: list[str] = []
    for b in blocks:
        text = str(b.get("text") or "").strip()
        if not text:
            continue
        if not title and text.startswith("# "):
            title = text.lstrip("# ").strip()
            continue
        paras.append(text.lstrip("# ").strip() if text.startswith("# ") else text)
    return (title or "The Pack's brief"), paras


def _render_markdown(title: str, paras: list[str]) -> bytes:
    return (f"# {title}\n\n" + "\n\n".join(paras)).encode("utf-8")


def _render_html(title: str, paras: list[str]) -> bytes:
    body = _markdown.markdown("\n\n".join(paras))
    html = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<title>{escape(title)}</title>"
        "<style>body{font:16px/1.6 -apple-system,Segoe UI,sans-serif;max-width:720px;"
        "margin:40px auto;padding:0 16px;color:#1a1a1a}h1{font-size:28px}</style></head>"
        f"<body><h1>{escape(title)}</h1>{body}</body></html>"
    )
    return html.encode("utf-8")


def _render_pdf(title: str, paras: list[str]) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=LETTER, title=title)
    styles = getSampleStyleSheet()
    story = [Paragraph(escape(title), styles["Title"]), Spacer(1, 14)]
    for p in paras:
        story.append(Paragraph(escape(p), styles["BodyText"]))
        story.append(Spacer(1, 8))
    doc.build(story)
    return buf.getvalue()


def _render_docx(title: str, paras: list[str]) -> bytes:
    d = Document()
    d.add_heading(title, level=0)
    for p in paras:
        d.add_paragraph(p)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _render_xlsx(title: str, paras: list[str]) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Brief"
    ws["A1"] = title
    ws["A3"], ws["B3"] = "#", "Finding"
    for i, p in enumerate(paras, start=1):
        ws.cell(row=3 + i, column=1, value=i)
        ws.cell(row=3 + i, column=2, value=p)
    ws.column_dimensions["B"].width = 100
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _render_pptx(title: str, paras: list[str]) -> bytes:
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title
    for i, p in enumerate(paras, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Finding {i}"
        slide.placeholders[1].text = p
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _render_png(title: str, paras: list[str]) -> bytes:
    img = Image.new("RGB", (1200, 630), (15, 15, 15))
    draw = ImageDraw.Draw(img)
    try:
        big = ImageFont.truetype("arial.ttf", 44)
        small = ImageFont.truetype("arial.ttf", 26)
    except Exception:  # noqa: BLE001 — no system font: fall back to the bundled bitmap font
        big = small = ImageFont.load_default()
    draw.text((60, 56), title[:60], fill=(255, 255, 255), font=big)
    y = 150
    for p in paras[:6]:
        draw.text((60, y), ("- " + p)[:96], fill=(180, 180, 188), font=small)
        y += 48
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


_RENDERERS = {
    "md": _render_markdown,
    "html": _render_html,
    "pdf": _render_pdf,
    "docx": _render_docx,
    "xlsx": _render_xlsx,
    "pptx": _render_pptx,
    "png": _render_png,
}


def forge(
    blocks: list[dict], title: str = "", formats: list[str] | None = None
) -> dict[str, bytes]:
    """Render the brief's blocks to the requested document formats. Returns {format: bytes}. A
    renderer that raises is skipped (best-effort) rather than sinking the whole Forge."""
    title, paras = _title_and_paras(blocks, title)
    out: dict[str, bytes] = {}
    for fmt in formats or list(_RENDERERS):
        renderer = _RENDERERS.get(fmt)
        if renderer is None:
            continue
        try:
            out[fmt] = renderer(title, paras)
        except Exception:  # noqa: BLE001 — one bad format never sinks the rest
            continue
    return out
